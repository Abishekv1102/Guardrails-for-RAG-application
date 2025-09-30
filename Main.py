import streamlit as st
import os
import tempfile
import google.generativeai as genai
from PyPDF2 import PdfReader
from langchain_core.documents import Document
from pptx import Presentation
from langchain.text_splitter import CharacterTextSplitter
from langchain_chroma import Chroma
from langchain_google_genai import GoogleGenerativeAIEmbeddings,ChatGoogleGenerativeAI
from langchain.chains import create_retrieval_chain
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_core.prompts import ChatPromptTemplate
from docx import Document as DocxDocument
import guardrails_1 as app
from dotenv import load_dotenv

#Load environment variable and API Key
load_dotenv()
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
gemini_model = genai.GenerativeModel("gemini-1.5-flash")


# File Processing Functions
def process_file(file):
   """Handle different file formats"""
   with tempfile.NamedTemporaryFile(delete=False) as temp:
      temp.write(file.read())
      temp_path = temp.name
   
   try:
      if file.type == "text/plain":
         return process_txt(temp_path)
      elif file.type == "application/pdf":
         return process_pdf(temp_path)
      elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
         return process_docx(temp_path)
      elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
         return process_pptx(temp_path)
   finally:
      os.unlink(temp_path)
def process_txt(path):
   all_text = ""
   for file in uploaded_files:
      file.seek(0)
      text = file.read().decode("utf-8")
      if text.strip():
         all_text += text + "\n"
   return all_text
def process_pdf(path):
   return "\n".join([page.extract_text() for page in PdfReader(path).pages])
def process_docx(path):
   doc = DocxDocument(path)
   return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
def process_pptx(path):
   prs = Presentation(path)
   text_runs = []
   for slide in prs.slides:
      for shape in slide.shapes:
         if hasattr(shape, "text"):
            text = shape.text.strip()
            if text:
               text_runs.append(text)
   return "\n".join(text_runs)
# Streamlit App
st.title("🔐 RAG Application with Guardrails")
st.sidebar.header("📂 Upload Your Documents")
# File Upload
uploaded_files = st.sidebar.file_uploader(
   "Upload documents (TXT/PDF/DOCX/PPTX)",
   type=["txt", "pdf", "docx", "pptx"],
   accept_multiple_files=True
)
# Document Processing
if uploaded_files:
   # Process files
   all_text = "\n".join([process_file(file) for file in uploaded_files])
   
   # Split text
   text_splitter = CharacterTextSplitter(
      separator="\n",
      chunk_size=1000,
      chunk_overlap=200
   )
   chunks = text_splitter.split_text(all_text)
   documents = [Document(page_content=chunk) for chunk in chunks]
   
   # Create Vector Store with Google Embeddings
   vectorstore = Chroma.from_documents(documents=documents, embedding=GoogleGenerativeAIEmbeddings(model="models/embedding-001",google_api_key=GEMINI_API_KEY,transport="rest", request_timeout=120))
   
   # Q&A Interface
   st.header("💬 Ask About Your Documents")
   question = st.text_input("Enter your question:")
   
   if question:
      retriever = vectorstore.as_retriever(search_type="similarity", search_kwargs={"k": 3})
      retrieved_docs = retriever.invoke(question)
      # Retrieve Context
      context = "\n\n".join([doc.page_content for doc in retrieved_docs])
      #LLM model
      llm = ChatGoogleGenerativeAI(model="gemini-1.5-flash",temperature=0.3, max_tokens=500,google_api_key=GEMINI_API_KEY)
      system_prompt = (
         "You are an assistant for question-answering tasks. "
         "Use the following pieces of retrieved context to answer "
         "the question. If you don't know the answer, say that you "
         "don't know. Use three sentences maximum and keep the "
         "answer concise."
         "Return your answer as a JSON object with the key 'RAG application response'."
         "\n\n"
         "{context}"
      )
      prompt = ChatPromptTemplate.from_messages(
         [
            ("system", system_prompt),
            ("human", "{input}"),
         ]
      )
      question_answer_chain = create_stuff_documents_chain(llm, prompt)
      rag_chain = create_retrieval_chain(retriever, question_answer_chain)
      response1 = rag_chain.invoke({"input": question})
      rag_generated_answer=response1["answer"]

      #Import guardrails
      guard=app.guard
      validation_outcome = guard.validate(rag_generated_answer)
      # Display validation results
      if validation_outcome.validation_passed:
         st.success("RAG response passed all guardrail checks!")
         st.write(validation_outcome.validated_output)
      else:
         st.error("RAG response failed guardrail checks!")
         st.write("Original RAG response:", rag_generated_answer)

